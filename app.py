import openai
import streamlit as st
from docx import Document
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor, as_completed
import time
import streamlit.components.v1 as components

# ----- CONFIGURATION -----
OPENAI_API_KEY = 'your-openai-api-key-here'
MODEL = 'gpt-4o'
MAX_TOKENS_PER_CHUNK = 2000
SUMMARY_MAX_TOKENS = 800

# ----- UTILITY FUNCTIONS -----
def read_docx(uploaded_file):
    doc = Document(uploaded_file)
    return [para.text.strip() for para in doc.paragraphs if para.text.strip()]

def read_pptx(uploaded_file):
    prs = Presentation(uploaded_file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    text_runs.append(text)
    return text_runs

def chunk_paragraphs(paragraphs, max_tokens=MAX_TOKENS_PER_CHUNK):
    chunks, current_chunk, current_tokens = [], [], 0
    for para in paragraphs:
        tokens = len(para.split())
        if current_tokens + tokens > max_tokens:
            chunks.append("\n".join(current_chunk))
            current_chunk, current_tokens = [para], tokens
        else:
            current_chunk.append(para)
            current_tokens += tokens
    if current_chunk:
        chunks.append("\n".join(current_chunk))
    return chunks

def summarize_chunk_safe(chunk, retry_count=3):
    for attempt in range(retry_count):
        try:
            response = openai.ChatCompletion.create(
                model=MODEL,
                api_key=OPENAI_API_KEY,
                messages=[
                    {"role": "system", "content": "Summarize the following document chunk clearly, retaining important requirements, features, and key points."},
                    {"role": "user", "content": chunk}
                ],
                temperature=0.2,
                max_tokens=1500
            )
            return response.chocies[0].message.content
        except Exception as e:
            print(f"Error summarizing chunk (attempt {attempt+1}): {e}")
            time.sleep(2)
    return "[Error: Failed to summarize this chunk.]"

def summarize_document(paragraphs):
    chunks = chunk_paragraphs(paragraphs)
    summaries = [""] * len(chunks)
    progress_bar = st.progress(0)
    total = len(chunks)

    with ThreadPoolExecutor(max_workers=5) as executor:
        futures = {executor.submit(summarize_chunk_safe, chunk): i for i, chunk in enumerate(chunks)}
        completed = 0
        for future in as_completed(futures):
            i = futures[future]
            try:
                summaries[i] = future.result()
            except Exception as e:
                summaries[i] = "[Error]"
            completed += 1
            progress_bar.progress(completed / total)

    progress_bar.empty()
    return "\n\n".join(summaries)

def generate_new_frd(existing_brd_summary, existing_frd_summary, new_brd_summary):
    user_prompt = f"""
EXISTING BRD SUMMARY:
{existing_brd_summary}

EXISTING FRD SUMMARY:
{existing_frd_summary}

NEW BRD SUMMARY:
{new_brd_summary}

Please generate the NEW FRD.
"""
    system_prompt = (
        "You are an expert business analyst. "
        "You are given summarized versions of an existing BRD, FRD, and a new BRD. "
        "Your task is to create a NEW FRD based on the new BRD, maintaining structure and clarity of the existing FRD."
    )
    response = openai.ChatCompletion.create(
        model=MODEL,
        api_key=OPENAI_API_KEY,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.2,
        max_tokens=3000
    )
    return response.choices[0].message.content

# ----- STREAMLIT UI -----
st.set_page_config(
    page_title="Business Analysis Toolkit",
    layout="wide",
    page_icon="🚀"
)

# Custom CSS for animations and styling
st.markdown("""
<style>
    /* Layout adjustments */
    .main .block-container {
        padding-top: 1rem;
    }
    .stApp {
        margin-top: 0;
    }
    .fade-in h1:first-child {
        margin-top: 0;
    }
    
    /* Animations */
    @keyframes fadeIn {
        from { opacity: 0; transform: translateY(20px); }
        to { opacity: 1; transform: translateY(0); }
    }
    
    @keyframes textGradient {
        0% { background-position: 0% 50%; }
        50% { background-position: 100% 50%; }
        100% { background-position: 0% 50%; }
    }
    
    @keyframes pulse {
        0% { transform: scale(1); }
        50% { transform: scale(1.05); }
        100% { transform: scale(1); }
    }
    
    .fade-in {
        animation: fadeIn 0.5s ease-out forwards;
    }
    
    .feature-header {
        background: linear-gradient(90deg, #4a6bff, #6a4bff, #4a6bff);
        background-size: 200% 200%;
        animation: textGradient 3s ease infinite, pulse 2s ease infinite;
        -webkit-background-clip: text;
        background-clip: text;
        color: transparent;
        display: inline-block;
        margin-bottom: 15px;
    }
    
    .sidebar .sidebar-content {
        background: linear-gradient(180deg, #4a6bff, #3a52d3);
        color: white;
    }
    
    .sidebar .sidebar-content .block-container {
        color: white;
    }
    
    .sidebar .sidebar-content .stRadio > div {
        color: white;
    }
    
    .sidebar .sidebar-content .stRadio > label > div:first-child {
        background-color: rgba(255,255,255,0.1);
        padding: 8px;
        border-radius: 8px;
        margin-bottom: 8px;
        transition: all 0.3s;
    }
    
    .sidebar .sidebar-content .stRadio > label > div:first-child:hover {
        background-color: rgba(255,255,255,0.2);
        transform: translateX(5px);
    }
    
    .stButton>button {
        background: linear-gradient(90deg, #4a6bff, #3a52d3);
        color: white;
        border: none;
        border-radius: 8px;
        padding: 10px 24px;
        font-weight: bold;
        transition: all 0.3s;
    }
    
    .stButton>button:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(0,0,0,0.1);
    }
    
    .stFileUploader>div>div>div>div {
        border: 2px dashed #4a6bff;
        border-radius: 12px;
        padding: 20px;
        background: rgba(74, 107, 255, 0.05);
    }
    
    .coming-soon {
        background: linear-gradient(45deg, #ff6b6b, #ff8e8e);
        color: white;
        padding: 20px;
        border-radius: 12px;
        text-align: center;
        margin-top: 20px;
        animation: pulse 2s ease infinite;
    }
    
    .feature-card {
        background: white;
        border-radius: 12px;
        padding: 20px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.1);
        margin-bottom: 20px;
        border-left: 5px solid #4a6bff;
    }
    
    .success-box {
        background: linear-gradient(45deg, #4CAF50, #8BC34A);
        color: white;
        padding: 20px;
        border-radius: 12px;
        margin-bottom: 20px;
    }
    
    .nav-item-selected {
        background-color: rgba(255,255,255,0.3) !important;
        font-weight: bold !important;
    }
    
    /* Sidebar animations */
    @keyframes sidebarFadeIn {
        from { opacity: 0; transform: translateX(-20px); }
        to { opacity: 1; transform: translateX(0); }
    }
    
    .sidebar .fade-in {
        animation: sidebarFadeIn 0.5s ease-out forwards;
    }
    
    .radio-option {
        padding: 12px;
        margin: 8px 0;
        border-radius: 8px;
        transition: all 0.3s ease;
        background-color: rgba(240, 242, 246, 0.7);
    }
    
    .radio-option:hover {
        background-color: rgba(200, 220, 255, 0.9);
        transform: translateX(5px);
    }
    
    .selected-option {
        background-color: #4a8cff !important;
        color: white !important;
        font-weight: bold;
        box-shadow: 0 4px 8px rgba(0,0,0,0.1);
    }
    
    .title-container {
        position: relative;
    }
    
    .title-container::after {
        content: "";
        position: absolute;
        bottom: -10px;
        left: 0;
        width: 100%;
        height: 3px;
        background: linear-gradient(90deg, #4a8cff, #ff6b6b, #4a8cff);
        background-size: 200% 200%;
        animation: gradient 3s ease infinite;
        border-radius: 3px;
    }
    
    @keyframes gradient {
        0% { background-position: 0% 50%; }
        50% { background-position: 100% 50%; }
        100% { background-position: 0% 50%; }
    }
</style>
""", unsafe_allow_html=True)

# JavaScript for additional interactivity
components.html("""
<script>
    document.addEventListener('DOMContentLoaded', function() {
        // Add click animation to radio options
        const radioOptions = document.querySelectorAll('.radio-option');
        radioOptions.forEach(option => {
            option.addEventListener('click', function() {
                this.style.transform = 'scale(0.95)';
                setTimeout(() => {
                    this.style.transform = '';
                }, 200);
            });
        });
    });
</script>
""")

# Sidebar content
with st.sidebar:
    st.markdown("""
    <div class="fade-in title-container">
        <h1 style="color: black; text-align: center; margin-bottom: 20px;" class="pulse-on-hover">GETTS</h1>
    </div>
    """, unsafe_allow_html=True)
    
    # Custom radio buttons with enhanced visuals
    options = ["Generate FRD", "Generate Test Scenario", "Generate Mockup", "Generate Excel File"]
    selected_topic = st.radio(
        "Select Functionality",
        options,
        index=0,
        key="nav",
        label_visibility="collapsed"
    )
    
    # Add custom styling to the radio buttons
    st.markdown(f"""
    <script>
        document.addEventListener('DOMContentLoaded', function() {{
            const options = {options};
            const radioButtons = parent.document.querySelectorAll('div[role="radiogroup"] label');
            
            radioButtons.forEach((button, index) => {{
                // Wrap each option in our styled div
                const wrapper = document.createElement('div');
                wrapper.className = 'radio-option';
                if(options[index] === "{selected_topic}") {{
                    wrapper.classList.add('selected-option');
                }}
                
                // Move the content into the wrapper
                while(button.firstChild) {{
                    wrapper.appendChild(button.firstChild);
                }}
                button.appendChild(wrapper);
                
                // Add click handler to update styling
                wrapper.addEventListener('click', function() {{
                    // Remove selected class from all options
                    document.querySelectorAll('.radio-option').forEach(opt => {{
                        opt.classList.remove('selected-option');
                    }});
                    // Add to clicked option
                    this.classList.add('selected-option');
                }});
            }});
        }});
    </script>
    """, unsafe_allow_html=True)

# Main content area - starts from top
st.markdown("""
<div class="fade-in">
    <h1 style="color: #4a6bff; margin-top: 0;">🚀 Business Analysis Toolkit</h1>
    <p style="color: #666; font-size: 1.1rem;">Streamline your business analysis workflow with powerful automation tools</p>
    <hr>
</div>
""", unsafe_allow_html=True)

if selected_topic == "Generate FRD":
    st.markdown(f"""
    <div class="feature-card fade-in">
        <h2 class="feature-header">📄 Auto FRD Generator</h2>
        <p>Upload your existing BRD, FRD, and (optionally) new BRD file below to generate a new FRD.</p>
    </div>
    """, unsafe_allow_html=True)
    
    col1, col2 = st.columns(2)
    
    with col1:
        existing_brd_file = st.file_uploader("Upload Existing BRD (.docx or .pptx)", type=["docx", "pptx"], key="brd")
        existing_frd_file = st.file_uploader("Upload Existing FRD (.docx)", type="docx", key="frd")
    
    with col2:
        new_brd_file = st.file_uploader("Upload New BRD (.docx)", type="docx", key="new_brd")
    
    if st.button("✨ Generate New FRD", type="primary"):
        if not existing_brd_file or not existing_frd_file:
            st.error("❌ Please upload both Existing BRD and Existing FRD.")
        else:
            with st.spinner("Reading and summarizing documents..."):
                if existing_brd_file.name.endswith(".pptx"):
                    paragraphs_brd = read_pptx(existing_brd_file)
                else:
                    paragraphs_brd = read_docx(existing_brd_file)

                paragraphs_frd = read_docx(existing_frd_file)
                paragraphs_new_brd = read_docx(new_brd_file) if new_brd_file else []

                summary_brd = summarize_document(paragraphs_brd)
                summary_frd = summarize_document(paragraphs_frd)
                summary_new_brd = summarize_document(paragraphs_new_brd) if new_brd_file else "No new BRD provided."

            with st.spinner("Generating NEW FRD..."):
                new_frd_text = generate_new_frd(summary_brd, summary_frd, summary_new_brd)

            st.markdown("""
            <div class="success-box fade-in">
                <h3 style="color: white; margin: 0;">✅ FRD Generated Successfully!</h3>
            </div>
            """, unsafe_allow_html=True)
            
            col1, col2 = st.columns([1, 3])
            with col1:
                st.download_button(
                    "📥 Download New FRD", 
                    new_frd_text, 
                    file_name="new_frd.txt", 
                    mime="text/plain",
                    type="primary"
                )
            with col2:
                st.text_area("Preview of Generated FRD", new_frd_text, height=300, label_visibility="collapsed")

elif selected_topic == "Generate Test Scenario":
    st.markdown(f"""
    <div class="feature-card fade-in">
        <h2 class="feature-header">🧪 Generate Test Scenario</h2>
        <p>Automatically create comprehensive test scenarios from your requirements documents.</p>
    </div>
    """, unsafe_allow_html=True)
    
    st.markdown("""
    <div class="coming-soon fade-in">
        <h3>🚧 Coming Soon!</h3>
        <p>We're working hard to bring you this feature in our next release.</p>
    </div>
    """, unsafe_allow_html=True)
    
elif selected_topic == "Generate Mockup":
    st.markdown(f"""
    <div class="feature-card fade-in">
        <h2 class="feature-header">🎨 Generate Mockup</h2>
        <p>Transform your requirements into visual mockups automatically.</p>
    </div>
    """, unsafe_allow_html=True)
    
    st.markdown("""
    <div class="coming-soon fade-in">
        <h3>🚧 Coming Soon!</h3>
        <p>We're working hard to bring you this feature in our next release.</p>
    </div>
    """, unsafe_allow_html=True)
    
elif selected_topic == "Generate Excel File":
    st.markdown(f"""
    <div class="feature-card fade-in">
        <h2 class="feature-header">📊 Generate Excel File</h2>
        <p>Create detailed Excel reports from your business documents.</p>
    </div>
    """, unsafe_allow_html=True)
    
    st.markdown("""
    <div class="coming-soon fade-in">
        <h3>🚧 Coming Soon!</h3>
        <p>We're working hard to bring you this feature in our next release.</p>
    </div>
    """, unsafe_allow_html=True)

# Footer with animation
st.markdown("""
<div class="fade-in">
    <hr>
    <div style="text-align: center; color: #666; padding: 20px;">
        <p>© 2023 Business Analysis Toolkit | Enterprise Edition</p>
    </div>
</div>
""", unsafe_allow_html=True)

# Add some confetti animation for success
if 'show_confetti' in st.session_state and st.session_state.show_confetti:
    components.html("""
    <script src="https://cdn.jsdelivr.net/npm/canvas-confetti@1.5.1/dist/confetti.browser.min.js"></script>
    <script>
        confetti({
            particleCount: 100,
            spread: 70,
            origin: { y: 0.6 }
        });
        setTimeout(() => {
            confetti({
                particleCount: 50,
                spread: 100,
                origin: { y: 0.6 }
            });
        }, 300);
    </script>
    """)
    st.session_state.show_confetti = False
