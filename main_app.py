# main_app.py

import streamlit as st
import docx
from pptx import Presentation
import concurrent.futures
import os
from langgraph_workflow import build_frd_graph
from openai import OpenAIError
import openai

# Set your OpenAI API key
openai.api_key = os.getenv("OPENAI_API_KEY")

# Utility: Read docx
def read_docx(file):
    doc = docx.Document(file)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

# Utility: Read pptx
def read_pptx(file):
    prs = Presentation(file)
    return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())

# Split into chunks for GPT
def chunk_paragraphs(text, max_tokens=1500):
    paragraphs = text.split("\n")
    chunks, current_chunk = [], ""
    for para in paragraphs:
        if len(current_chunk + para) < max_tokens:
            current_chunk += para + "\n"
        else:
            chunks.append(current_chunk.strip())
            current_chunk = para + "\n"
    if current_chunk:
        chunks.append(current_chunk.strip())
    return chunks

# Safe GPT summarizer
def summarize_chunk_safe(chunk):
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4-turbo",
            messages=[
                {"role": "system", "content": "Summarize in a business analyst style."},
                {"role": "user", "content": chunk}
            ],
            temperature=0.2
        )
        return response.choices[0].message.content.strip()
    except OpenAIError as e:
        return f"[Error summarizing chunk: {e}]"

# Parallel summarizer
def summarize_document(text):
    chunks = chunk_paragraphs(text)
    with concurrent.futures.ThreadPoolExecutor() as executor:
        summaries = list(executor.map(summarize_chunk_safe, chunks))
    return "\n".join(summaries)

# Streamlit UI
st.set_page_config(layout="wide", page_title="AI FRD Generator")
st.title("📄 AI-Powered FRD Generator with LangGraph")

option = st.sidebar.selectbox("Choose Functionality", ["Generate FRD"])

if option == "Generate FRD":
    st.subheader("Upload Documents")

    col1, col2 = st.columns(2)
    with col1:
        existing_brd_file = st.file_uploader("Upload Existing BRD", type=["docx", "pptx"])
        existing_frd_file = st.file_uploader("Upload Existing FRD", type=["docx", "pptx"])
    with col2:
        new_brd_file = st.file_uploader("Upload New BRD", type=["docx", "pptx"])
        user_notes = st.text_area("Additional Notes (Optional)", height=150)

    if st.button("Generate New FRD", type="primary"):
        if not all([existing_brd_file, existing_frd_file, new_brd_file]):
            st.error("Please upload all required documents.")
        else:
            with st.spinner("Reading and summarizing documents..."):
                def read_file(f): return read_docx(f) if f.name.endswith(".docx") else read_pptx(f)
                existing_brd_text = read_file(existing_brd_file)
                existing_frd_text = read_file(existing_frd_file)
                new_brd_text = read_file(new_brd_file)

                summary_brd = summarize_document(existing_brd_text)
                summary_frd = summarize_document(existing_frd_text)
                summary_new_brd = summarize_document(new_brd_text)

            with st.spinner("Generating FRD using LangGraph..."):
                try:
                    graph = build_frd_graph()
                    result = graph.invoke({
                        "existing_brd": summary_brd,
                        "existing_frd": summary_frd,
                        "new_brd": summary_new_brd,
                        "user_notes": user_notes
                    })
                    new_frd_text = result["new_frd"]
                    st.success("✅ FRD Generated Successfully!")
                    st.download_button("Download New FRD (txt)", new_frd_text, file_name="Generated_FRD.txt")

                except Exception as e:
                    st.error(f"Failed to generate FRD: {e}")
